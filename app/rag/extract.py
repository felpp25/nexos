"""Extracao de texto de PDF, DOCX, PPTX, imagens, TXT/MD/CSV.

Cada extractor devolve uma lista de blocos (texto, localizacao) para que a
citacao possa apontar "pagina 3" ou "slide 7".
"""
from __future__ import annotations

from pathlib import Path
from typing import Callable

Block = tuple[str, str]  # (texto, localizacao)

PDF_EXT = {".pdf"}
DOC_EXT = {".docx"}
SLIDE_EXT = {".pptx"}
IMG_EXT = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".gif"}
TEXT_EXT = {".txt", ".md", ".markdown", ".csv", ".json", ".log"}

SUPPORTED = PDF_EXT | DOC_EXT | SLIDE_EXT | IMG_EXT | TEXT_EXT


class ExtractionError(RuntimeError):
    pass


def kind_of(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext in PDF_EXT:
        return "pdf"
    if ext in DOC_EXT:
        return "docx"
    if ext in SLIDE_EXT:
        return "pptx"
    if ext in IMG_EXT:
        return "image"
    if ext in TEXT_EXT:
        return "text"
    return "other"


def _extract_pdf(path: Path) -> list[Block]:
    try:
        import pymupdf as fitz
    except ImportError:
        try:
            import fitz  # PyMuPDF < 1.24
        except ImportError as exc:
            raise ExtractionError("PyMuPDF nao instalado (pip install pymupdf)") from exc

    blocks: list[Block] = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            text = (page.get_text("text") or "").strip()
            if text:
                blocks.append((text, f"pagina {i}"))
    if not blocks:
        raise ExtractionError(
            "PDF sem texto selecionavel (provavelmente digitalizado). "
            "Converta para PDF pesquisavel ou envie o texto em .docx/.txt."
        )
    return blocks


def _extract_docx(path: Path) -> list[Block]:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise ExtractionError("python-docx nao instalado (pip install python-docx)") from exc

    document = docx.Document(str(path))
    parts: list[str] = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    if not parts:
        raise ExtractionError("Nenhum texto encontrado no .docx")
    return [("\n".join(parts), "documento")]


def _extract_pptx(path: Path) -> list[Block]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ExtractionError("python-pptx nao instalado (pip install python-pptx)") from exc

    prs = Presentation(str(path))
    blocks: list[Block] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[notas do apresentador] {notes}")
        if parts:
            blocks.append(("\n".join(parts), f"slide {i}"))
    if not blocks:
        raise ExtractionError("Nenhum texto encontrado no .pptx")
    return blocks


def _extract_image(path: Path) -> list[Block]:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise ExtractionError(
            "OCR indisponivel. Instale: pip install pytesseract pillow "
            "e o Tesseract OCR no Windows (winget install UB-Mannheim.TesseractOCR)."
        ) from exc

    try:
        text = pytesseract.image_to_string(Image.open(path), lang="por+eng").strip()
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"Falha no OCR ({exc}). Confirme que o Tesseract esta instalado e no PATH."
        ) from exc
    if not text:
        raise ExtractionError("OCR nao encontrou texto legivel nesta imagem.")
    return [(text, "imagem")]


def _extract_text(path: Path) -> list[Block]:
    raw = path.read_bytes()
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            text = raw.decode(enc).strip()
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ExtractionError("Nao foi possivel decodificar o arquivo de texto.")
    if not text:
        raise ExtractionError("Arquivo de texto vazio.")
    return [(text, "arquivo")]


_EXTRACTORS: dict[str, Callable[[Path], list[Block]]] = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "image": _extract_image,
    "text": _extract_text,
}


def extract(path: Path, filename: str | None = None) -> list[Block]:
    kind = kind_of(filename or path.name)
    extractor = _EXTRACTORS.get(kind)
    if extractor is None:
        raise ExtractionError(
            f"Formato nao suportado: {Path(filename or path.name).suffix or 'desconhecido'}. "
            "Aceitos: PDF, DOCX, PPTX, imagens (OCR), TXT, MD, CSV."
        )
    return extractor(path)
